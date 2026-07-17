"""Generic shape helpers. Every helper sets shape.name for testability."""
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from deck import theme


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def paint_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _style_runs(tf, size, color, bold, italic, align, line_spacing, space_after):
    for para in tf.paragraphs:
        para.alignment = align
        if line_spacing:
            para.line_spacing = line_spacing
        if space_after:
            para.space_after = space_after
        for run in para.runs or [para.add_run()]:
            run.font.name = theme.FONT
            run.font.size = size
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color


def tb(slide, x, y, w, h, text, size=None, color=None, bold=False,
       align=PP_ALIGN.LEFT, name=None, italic=False, line_spacing=None,
       space_after=None):
    shape = slide.shapes.add_textbox(x, y, w, h)
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    tf.text = lines[0]
    for line in lines[1:]:
        tf.add_paragraph().text = line
    _style_runs(tf, size or theme.BODY_SIZE, color or theme.DARK_TEXT,
                bold, italic, align, line_spacing, space_after)
    return shape


def box(slide, x, y, w, h, text, fill, text_color=None, size=None, bold=False,
        name="box", shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, italic=False,
        align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shape.adjustments[0] = 0.10
        except (IndexError, ValueError):
            pass
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n")
    tf.text = lines[0]
    for line in lines[1:]:
        tf.add_paragraph().text = line
    _style_runs(tf, size or theme.SMALL_SIZE, text_color or theme.WHITE,
                bold, italic, align, None, None)
    return shape


def panel(slide, x, y, w, h, fill=None, name="panel"):
    """Background container card; content is drawn on top of it by design."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or theme.ICE
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.06
    except (IndexError, ValueError):
        pass
    return shape


def add_title(slide, spec, color=None):
    tb(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(1.15),
       spec["title"], size=theme.TITLE_SIZE, color=color or theme.DEEP_BLUE,
       bold=True, name="title")


def add_killer(slide, text, dark=False):
    if dark:
        return tb(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.75),
                  text, size=theme.KILLER_SIZE, color=theme.WHITE, bold=True,
                  italic=True, name="killer")
    shape = box(slide, Inches(0.5), Inches(6.12), Inches(12.3), Inches(0.78),
                text, theme.ICE, text_color=theme.DEEP_BLUE,
                size=theme.KILLER_SIZE, bold=True, italic=True, name="killer")
    return shape


def add_bullets(slide, items, x, y, w, h, size=None, color=None):
    return tb(slide, x, y, w, h, "\n".join("•  " + i for i in items),
              size=size or theme.BODY_SIZE, color=color or theme.SLATE,
              name="bullets", line_spacing=1.12, space_after=Pt(5))


def add_footer(slide, spec, dark=False):
    text = f"{theme.PRODUCT_NAME}  ·  {spec['n']}/18"
    if spec["disclaimer"]:
        text += f"\n{theme.DISCLAIMER_VN}  ({theme.DISCLAIMER_EN})"
    tb(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.5),
       text, size=theme.FOOTER_SIZE,
       color=theme.ICE_TEXT if dark else theme.GRAY, name="footer")


def add_placeholder(slide, x, y, w, h, label):
    shape = box(slide, x, y, w, h, label, theme.LIGHT_GRAY,
                text_color=theme.GRAY, name="screenshot_placeholder",
                shape_type=MSO_SHAPE.RECTANGLE)
    return shape
