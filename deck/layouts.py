"""Layout renderers. This task: everything renders as 'standard'.
Later tasks replace entries in RENDERERS with specialized functions."""
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from deck import builders, content, theme


def render_standard(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    if spec.get("bullets"):
        builders.add_bullets(slide, spec["bullets"],
                             Inches(0.7), Inches(1.5), Inches(11.9), Inches(4.4))
    if spec.get("killer"):
        builders.add_killer(slide, spec["killer"])
    builders.add_footer(slide, spec)
    return slide


def render_hook(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    x = spec["extra"]
    builders.box(slide, Inches(0.7), Inches(2.0), Inches(5.4), Inches(2.2),
                 x["quote"], theme.DEEP_BLUE, size=theme.KILLER_SIZE, name="quote")
    builders.tb(slide, Inches(0.7), Inches(4.3), Inches(5.4), Inches(0.5),
                x["quote_by"], size=theme.SMALL_SIZE, color=theme.GRAY)
    for i, label in enumerate(x["doc_labels"]):
        row, col = divmod(i, 3)
        builders.box(slide,
                     Inches(6.8) + col * Inches(2.1),
                     Inches(1.7) + row * Inches(1.15),
                     Inches(1.95), Inches(1.0),
                     label, theme.LIGHT_GRAY, text_color=theme.DARK_TEXT,
                     name="doc_thumb")
    builders.add_footer(slide, spec)
    return slide


def render_product(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    x = spec["extra"]
    builders.tb(slide, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.6),
                x["pitch"], size=theme.KILLER_SIZE, color=theme.DARK_TEXT, bold=True)
    builders.box(slide, Inches(4.4), Inches(3.1), Inches(3.0), Inches(1.3),
                 x["hub"], theme.ORANGE, name="hub", bold=True)
    positions = [(1.0, 2.0), (4.4, 1.85), (7.8, 2.0),
                 (1.0, 4.6), (4.4, 4.75), (7.8, 4.6)]
    for (px, py), role in zip(positions, x["roles"]):
        builders.box(slide, Inches(px), Inches(py), Inches(3.0), Inches(0.85),
                     role, theme.DEEP_BLUE, name="role_box")
    builders.add_placeholder(slide, Inches(11.0), Inches(1.9), Inches(2.0),
                             Inches(3.6), x["screenshot"])
    for i, promise in enumerate(x["promises"]):
        builders.box(slide, Inches(0.5) + i * Inches(4.2), Inches(6.15),
                     Inches(4.0), Inches(0.75), promise, theme.DEEP_BLUE,
                     name="promise", bold=True)
    builders.add_footer(slide, spec)
    return slide


def render_before_after(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    x = spec["extra"]
    builders.box(slide, Inches(0.6), Inches(1.35), Inches(6.0), Inches(0.5),
                 "TRƯỚC", theme.GRAY, name="col_before", bold=True)
    builders.box(slide, Inches(6.85), Inches(1.35), Inches(6.0), Inches(0.5),
                 "SAU", theme.ORANGE, name="col_after", bold=True)
    builders.add_bullets(slide, x["before"], Inches(0.6), Inches(2.0),
                         Inches(6.0), Inches(3.2), size=theme.SMALL_SIZE)
    builders.add_bullets(slide, x["after"], Inches(6.85), Inches(2.0),
                         Inches(6.0), Inches(3.2), size=theme.SMALL_SIZE)
    builders.box(slide, Inches(0.6), Inches(5.35), Inches(6.0), Inches(0.5),
                 x["before_bar"], theme.RED, name="time_bar", bold=True)
    builders.box(slide, Inches(6.85), Inches(5.35), Inches(2.4), Inches(0.5),
                 x["after_bar"], theme.GREEN, name="time_bar", bold=True)
    builders.add_killer(slide, spec["killer"])
    builders.add_footer(slide, spec)
    return slide


def render_storyboard(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    x = spec["extra"]
    for i, step in enumerate(x["steps"]):
        row, col = divmod(i, 3)
        builders.box(slide,
                     Inches(0.6) + col * Inches(4.3),
                     Inches(1.5) + row * Inches(2.5),
                     Inches(4.1), Inches(2.3),
                     f"{i + 1}. {step}", theme.DEEP_BLUE, name="story_card")
    builders.add_placeholder(slide, Inches(0.6), Inches(6.4), Inches(12.3),
                             Inches(0.45), x["screenshot"])
    builders.add_footer(slide, spec)
    return slide


def render_curve(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    x = spec["extra"]
    heights = [Inches(1.0), Inches(1.6), Inches(2.2)]
    xs = [Inches(0.8), Inches(4.0), Inches(9.6)]
    ys = [Inches(4.6), Inches(4.0), Inches(3.4)]
    for i, stage in enumerate(x["stages"]):
        builders.box(slide, xs[i], ys[i], Inches(2.9), heights[i],
                     stage, theme.DEEP_BLUE, name="stage_box", bold=True)
    builders.box(slide, Inches(7.1), Inches(3.9), Inches(2.3), Inches(1.2),
                 x["gap"], theme.RED, name="gap_box", bold=True)
    builders.add_bullets(slide, spec["bullets"], Inches(0.8), Inches(1.4),
                         Inches(11.9), Inches(1.7), size=theme.SMALL_SIZE)
    builders.add_killer(slide, spec["killer"])
    builders.add_footer(slide, spec)
    return slide


def render_pipeline(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    for i, step in enumerate(spec["extra"]["steps"]):
        row, col = divmod(i, 4)
        builders.box(slide,
                     Inches(0.5) + col * Inches(3.15),
                     Inches(1.9) + row * Inches(1.9),
                     Inches(3.0), Inches(1.5),
                     step, theme.DEEP_BLUE if row == 0 else theme.ORANGE,
                     name="chevron", shape_type=MSO_SHAPE.CHEVRON)
    builders.add_killer(slide, spec["killer"])
    builders.add_footer(slide, spec)
    return slide


def render_provenance(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    x = spec["extra"]
    for i, item in enumerate(x["chain"]):
        builders.box(slide, Inches(0.6), Inches(1.35) + i * Inches(0.72),
                     Inches(5.6), Inches(0.6), item, theme.DEEP_BLUE,
                     name="chain_box")
    builders.tb(slide, Inches(6.7), Inches(1.35), Inches(6.1), Inches(0.4),
                "Cấu trúc một bản ghi khoảng trống:", size=theme.SMALL_SIZE,
                color=theme.DEEP_BLUE, bold=True)
    builders.add_bullets(slide, x["gap_panel"], Inches(6.7), Inches(1.8),
                         Inches(6.1), Inches(3.2), size=theme.SMALL_SIZE)
    builders.add_placeholder(slide, Inches(6.7), Inches(5.1), Inches(6.1),
                             Inches(1.3), x["screenshot"])
    builders.add_footer(slide, spec)
    return slide


def render_grounding(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    x = spec["extra"]
    for i, src in enumerate(x["sources"]):
        row, col = divmod(i, 2)
        builders.box(slide,
                     Inches(0.6) + col * Inches(3.0),
                     Inches(1.6) + row * Inches(1.15),
                     Inches(2.85), Inches(1.0),
                     src, theme.DEEP_BLUE, name="source_box")
    builders.box(slide, Inches(7.0), Inches(2.4), Inches(2.6), Inches(1.6),
                 x["layer"], theme.ORANGE, name="layer_box", bold=True)
    builders.box(slide, Inches(10.1), Inches(2.0), Inches(2.7), Inches(1.0),
                 x["out"], theme.GREEN, name="out_box")
    builders.box(slide, Inches(10.1), Inches(3.4), Inches(2.7), Inches(1.2),
                 x["abstain"], theme.RED, name="abstain_box")
    builders.add_bullets(slide, spec["bullets"], Inches(0.6), Inches(5.2),
                         Inches(12.1), Inches(0.9), size=theme.SMALL_SIZE)
    builders.add_footer(slide, spec)
    return slide


def render_architecture(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    x = spec["extra"]
    for i, band in enumerate(x["bands"]):
        builders.box(slide, Inches(0.6), Inches(1.4) + i * Inches(1.0),
                     Inches(8.6), Inches(0.85), band, theme.DEEP_BLUE,
                     name="band")
    builders.tb(slide, Inches(9.5), Inches(1.4), Inches(3.3), Inches(0.4),
                "Ranh giới tin cậy:", size=theme.SMALL_SIZE,
                color=theme.DEEP_BLUE, bold=True)
    builders.add_bullets(slide, x["trust"], Inches(9.5), Inches(1.85),
                         Inches(3.3), Inches(2.6), size=theme.SMALL_SIZE)
    builders.tb(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.4),
                x["note"], size=theme.SMALL_SIZE, color=theme.GRAY)
    builders.add_killer(slide, spec["killer"])
    builders.add_footer(slide, spec)
    return slide


def _style_table(frame, header_fill, body_size):
    table = frame.table
    for r, row in enumerate(table.rows):
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else theme.WHITE
            for para in cell.text_frame.paragraphs:
                for run in para.runs or [para.add_run()]:
                    run.font.name = theme.FONT
                    run.font.size = body_size
                    run.font.bold = r == 0
                    run.font.color.rgb = (theme.WHITE if r == 0
                                          else theme.DARK_TEXT)


def _add_table(slide, name, header, rows, x, y, w, h, body_size=None):
    frame = slide.shapes.add_table(len(rows) + 1, len(header), x, y, w, h)
    frame.name = name
    table = frame.table
    for c, text in enumerate(header):
        table.cell(0, c).text = text
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    _style_table(frame, theme.DEEP_BLUE, body_size or theme.SMALL_SIZE)
    return frame


def render_compare_table(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    x = spec["extra"]
    _add_table(slide, "compare_table", x["cols"], x["rows"],
               Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.9))
    builders.add_footer(slide, spec)
    return slide


def render_criteria(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    x = spec["extra"]
    _add_table(slide, "criteria_table",
               ["Đề bài yêu cầu", "Chúng tôi có trong demo"], x["rows"],
               Inches(0.5), Inches(1.35), Inches(12.3), Inches(3.9))
    for i, benefit in enumerate(x["benefits"]):
        builders.box(slide, Inches(0.5) + i * Inches(3.15), Inches(5.6),
                     Inches(3.0), Inches(1.0), benefit, theme.ORANGE,
                     name="benefit")
    builders.add_footer(slide, spec)
    return slide


def render_validation(prs, spec):
    slide = builders.add_blank(prs)
    builders.add_title(slide, spec)
    builders.add_bullets(slide, spec["bullets"], Inches(0.5), Inches(1.3),
                         Inches(12.3), Inches(1.1), size=theme.SMALL_SIZE)
    _add_table(slide, "metric_table",
               ["Chỉ số", "Kết quả", "Ghi chú"], spec["extra"]["metrics"],
               Inches(0.5), Inches(2.5), Inches(12.3), Inches(3.4))
    builders.add_killer(slide, spec["killer"])
    builders.add_footer(slide, spec)
    return slide


RENDERERS = {name: render_standard for name in content.LAYOUTS}
RENDERERS.update({
    "hook": render_hook,
    "product": render_product,
    "before_after": render_before_after,
    "storyboard": render_storyboard,
})
RENDERERS.update({
    "curve": render_curve,
    "pipeline": render_pipeline,
    "provenance": render_provenance,
    "grounding": render_grounding,
    "architecture": render_architecture,
})
RENDERERS.update({
    "compare_table": render_compare_table,
    "criteria": render_criteria,
    "validation": render_validation,
})
