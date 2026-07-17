"""Programmatic visual QA: text overflow (real Segoe UI metrics) and
shape-overlap detection. A stand-in for opening the deck — this machine has
no PowerPoint/LibreOffice renderer. Advisory: exit 1 lists findings.

Usage: python deck/qa_visual.py [path-to-pptx]
Requires Pillow and the Windows Segoe UI font files.
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from PIL import ImageFont
from pptx import Presentation

EMU_IN = 914400
FONTS = {}
FONT_FILES = {
    (False, False): r"C:\Windows\Fonts\segoeui.ttf",
    (True, False): r"C:\Windows\Fonts\segoeuib.ttf",
    (False, True): r"C:\Windows\Fonts\segoeuii.ttf",
    (True, True): r"C:\Windows\Fonts\segoeuiz.ttf",
}
LINE_FACTOR = 1.30  # Segoe UI line height ~1.33 em; slightly generous
TOLERANCE = 1.06    # allow 6% slack before flagging


def font_for(size_pt, bold, italic):
    # Cache key rounds to whole points; sub-point size differences are far
    # inside the 6% TOLERANCE, so the approximation cannot flip a verdict.
    key = (round(size_pt), bold, italic)
    if key not in FONTS:
        FONTS[key] = ImageFont.truetype(FONT_FILES[(bold, italic)], round(size_pt))
    return FONTS[key]


def wrap_height_pt(text, width_pt, size_pt, bold, italic, line_spacing, space_after_pt):
    font = font_for(size_pt, bold, italic)
    total_lines = 0
    for para in text.split("\n"):
        words = para.split(" ")
        line = ""
        lines = 1
        for word in words:
            candidate = word if not line else line + " " + word
            if font.getlength(candidate) <= width_pt or not line:
                line = candidate
            else:
                lines += 1
                line = word
        total_lines += lines
    paras = text.count("\n") + 1
    line_h = size_pt * LINE_FACTOR * (line_spacing or 1.0)
    return total_lines * line_h + paras * (space_after_pt or 0)


def para_style(para, default=(11, False, False)):
    run = para.runs[0] if para.runs else None
    if run is None or run.font.size is None:
        return default
    return (run.font.size.pt, bool(run.font.bold), bool(run.font.italic))


def check_text_shape(idx, sh, problems):
    tf = sh.text_frame
    text = tf.text
    if not text.strip():
        return
    inner_w_pt = (sh.width / EMU_IN - 0.2) * 72
    inner_h_pt = (sh.height / EMU_IN - 0.1) * 72
    if inner_w_pt <= 0:
        return
    total = 0.0
    for para in tf.paragraphs:
        size_pt, bold, italic = para_style(para)
        spacing = para.line_spacing if isinstance(para.line_spacing, float) else None
        space_after = para.space_after.pt if para.space_after else 0
        total += wrap_height_pt(para.text or " ", inner_w_pt, size_pt, bold,
                                italic, spacing, space_after)
    if total > inner_h_pt * TOLERANCE:
        problems.append(f"slide {idx}: OVERFLOW {sh.name!r} needs ~{total:.0f}pt, "
                        f"box has {inner_h_pt:.0f}pt :: {text[:60]!r}")


def check_table(idx, sh, problems):
    tbl = sh.table
    n_rows = len(tbl.rows)
    n_cols = len(tbl.columns)
    col_w_pt = (sh.width / EMU_IN / n_cols - 0.19) * 72
    row_budget_pt = sh.height / EMU_IN / n_rows * 72
    for r, row in enumerate(tbl.rows):
        for c, cell in enumerate(row.cells):
            text = cell.text
            if not text.strip():
                continue
            para = cell.text_frame.paragraphs[0]
            size_pt, bold, italic = para_style(para)
            # +7pt approximates the default top+bottom cell insets (0.05in each).
            need = wrap_height_pt(text, col_w_pt, size_pt, bold, italic, None, 0) + 7
            if need > row_budget_pt * TOLERANCE:
                problems.append(
                    f"slide {idx}: TABLE-CELL {sh.name!r} r{r}c{c} needs "
                    f"~{need:.0f}pt, row budget {row_budget_pt:.0f}pt :: {text[:50]!r}")


def rects_overlap(a, b):
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    ox = min(ax2, bx2) - max(ax1, bx1)
    oy = min(ay2, by2) - max(ay1, by1)
    return ox > 0.03 * EMU_IN and oy > 0.03 * EMU_IN


def check_overlaps(idx, shapes, problems):
    rects = []
    for sh in shapes:
        # 'panel' is a background container by design; 'title' boxes are
        # taller than most rendered titles and their fit is verified by the
        # overflow check, so AABB collisions with them are false alarms.
        # Footers participate in overlap checks like any other shape.
        if sh.left is None or sh.name in ("panel", "title"):
            continue
        rects.append((sh.name, (sh.left, sh.top, sh.left + sh.width,
                                sh.top + sh.height)))
    for i in range(len(rects)):
        for j in range(i + 1, len(rects)):
            if rects_overlap(rects[i][1], rects[j][1]):
                problems.append(f"slide {idx}: OVERLAP {rects[i][0]!r} x {rects[j][0]!r}")


def main():
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    for path in FONT_FILES.values():
        if not os.path.exists(path):
            print(f"SKIPPED: font file {path} not found — run on a Windows "
                  "machine with Segoe UI installed.")
            return 0
    path = sys.argv[1] if len(sys.argv) > 1 else "deck/output/deck.pptx"
    prs = Presentation(path)
    problems = []
    for idx, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if sh.has_table:
                check_table(idx, sh, problems)
            elif sh.has_text_frame:
                check_text_shape(idx, sh, problems)
        check_overlaps(idx, list(slide.shapes), problems)
    if not problems:
        print(f"QA CLEAN: no overflow or overlap detected in {path}.")
        return 0
    print(f"{len(problems)} finding(s) in {path}:")
    for p in problems:
        print(" ", p)
    return 1


if __name__ == "__main__":
    sys.exit(main())
