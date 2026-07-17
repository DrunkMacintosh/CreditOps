import pytest
from pptx import Presentation
from pptx.util import Inches

from deck.content import SLIDES


@pytest.fixture(scope="session")
def prs():
    from deck.build import build_deck
    return Presentation(build_deck("deck/output/test_deck.pptx"))


def slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return "\n".join(parts)


def test_sixteen_by_nine(prs):
    assert prs.slide_width == Inches(13.333)
    assert prs.slide_height == Inches(7.5)


def test_eighteen_slides(prs):
    assert len(list(prs.slides)) == 18


def test_every_title_present(prs):
    for spec, slide in zip(SLIDES, prs.slides):
        assert spec["title"] in slide_text(slide), f"slide {spec['n']} title missing"


def test_every_footer_has_product_name(prs):
    from deck.theme import PRODUCT_NAME
    for spec, slide in zip(SLIDES, prs.slides):
        assert PRODUCT_NAME in slide_text(slide), f"slide {spec['n']} footer missing"


def test_speaker_notes_written(prs):
    slides = list(prs.slides)
    assert "chuẩn bị" in slides[0].notes_slide.notes_text_frame.text


def names(slide):
    return [s.name for s in slide.shapes]


def test_hook_slide_has_doc_wall(prs):
    slide = list(prs.slides)[0]
    assert names(slide).count("doc_thumb") == 12
    assert "Minh An" in slide_text(slide)


def test_product_slide_roles_and_placeholder(prs):
    slide = list(prs.slides)[3]
    assert names(slide).count("role_box") == 6
    assert "hub" in names(slide)
    assert "screenshot_placeholder" in names(slide)


def test_before_after_columns(prs):
    slide = list(prs.slides)[4]
    ns = names(slide)
    assert "col_before" in ns and "col_after" in ns
    assert ns.count("time_bar") == 2
    assert "NHIỀU TUẦN" in slide_text(slide) and "VÀI NGÀY" in slide_text(slide)


def test_storyboard_six_cards(prs):
    slide = list(prs.slides)[5]
    assert names(slide).count("story_card") == 6
    assert "screenshot_placeholder" in names(slide)


def test_curve_slide(prs):
    slide = list(prs.slides)[2]
    ns = names(slide)
    assert ns.count("stage_box") == 3 and "gap_box" in ns


def test_pipeline_eight_chevrons(prs):
    slide = list(prs.slides)[6]
    assert names(slide).count("chevron") == 8
    assert "1. Tiếp nhận nhu cầu" in slide_text(slide)


def test_provenance_chain(prs):
    slide = list(prs.slides)[7]
    ns = names(slide)
    assert ns.count("chain_box") == 7
    assert "screenshot_placeholder" in ns


def test_grounding_flow(prs):
    slide = list(prs.slides)[8]
    ns = names(slide)
    assert ns.count("source_box") == 6
    assert "layer_box" in ns and "abstain_box" in ns
    assert "out_box" in ns


def test_architecture_bands(prs):
    slide = list(prs.slides)[9]
    assert names(slide).count("band") == 4
    assert "Qwen3-30B-A3B" in slide_text(slide)


def get_table(slide, name):
    for shape in slide.shapes:
        if shape.name == name and shape.has_table:
            return shape.table
    return None


def test_compare_table(prs):
    table = get_table(list(prs.slides)[10], "compare_table")
    assert table is not None
    assert len(table.rows) == 8 and len(table.columns) == 5
    assert table.cell(0, 4).text == "SHB CreditOps EvidenceGraph"


def test_criteria_table(prs):
    table = get_table(list(prs.slides)[11], "criteria_table")
    assert table is not None and len(table.rows) == 7
    assert len(table.columns) == 2
    assert "GenAI" in slide_text(list(prs.slides)[11])


def test_validation_metrics(prs):
    slide = list(prs.slides)[12]
    table = get_table(slide, "metric_table")
    assert table is not None and len(table.rows) == 7
    assert len(table.columns) == 3
    assert "[X%]" in slide_text(slide)


def test_axes_slide(prs):
    assert names(list(prs.slides)[13]).count("axis_box") == 3


def test_roadmap_milestones(prs):
    assert names(list(prs.slides)[15]).count("milestone") == 5


def test_team_cards(prs):
    slide = list(prs.slides)[16]
    assert names(slide).count("member_card") == 5
    assert "[Họ tên]" in slide_text(slide)


def test_closing_slide(prs):
    slide = list(prs.slides)[17]
    ns = names(slide)
    assert ns.count("cta_box") == 3 and "qr_placeholder" in ns


def test_all_shapes_within_canvas(prs):
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if sh.left is None:
                continue
            assert sh.left >= 0 and sh.top >= 0, f"slide {i}: {sh.name} off-canvas"
            assert sh.left + sh.width <= prs.slide_width, f"slide {i}: {sh.name} exceeds width"
            assert sh.top + sh.height <= prs.slide_height, f"slide {i}: {sh.name} exceeds height"
